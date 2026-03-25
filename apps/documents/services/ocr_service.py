"""OCR and document text extraction service."""

import logging
import os
from pathlib import Path
from typing import Optional

import pytesseract
from PIL import Image
from pptx import Presentation
from pypdf import PdfReader
from docx import Document as DocxDocument

logger = logging.getLogger(__name__)


def extract_text_from_pdf(path: str) -> str:
    """Extract raw text from a PDF file."""
    text_parts = []
    try:
        reader = PdfReader(path)
        for page in reader.pages:
            page_text = page.extract_text() or ""
            text_parts.append(page_text)
    except Exception as ex:
        logger.warning("Failed to parse PDF '%s': %s", path, ex)
    return "\n".join(text_parts)


def extract_text_from_docx(path: str) -> str:
    """Extract raw text from a DOCX file."""
    text_parts = []
    try:
        doc = DocxDocument(path)
        for para in doc.paragraphs:
            text_parts.append(para.text)
    except Exception as ex:
        logger.warning("Failed to parse DOCX '%s': %s", path, ex)
    return "\n".join(text_parts)


def extract_text_from_pptx(path: str) -> str:
    """Extract raw text from a PPTX file."""
    text_parts = []
    try:
        prs = Presentation(path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_parts.append(shape.text)
    except Exception as ex:
        logger.warning("Failed to parse PPTX '%s': %s", path, ex)
    return "\n".join(text_parts)


def extract_text_from_image(path: str, lang: Optional[str] = "eng") -> str:
    """Extract text from an image using Tesseract."""
    try:
        image = Image.open(path)
        text = pytesseract.image_to_string(image, lang=lang)
        return text
    except Exception as ex:
        logger.warning("Failed to OCR image '%s': %s", path, ex)
        return ""


def extract_text_from_file(path: str) -> str:
    """Choose the correct extractor based on file extension."""
    extension = Path(path).suffix.lower()
    if extension == ".pdf":
        return extract_text_from_pdf(path)
    if extension == ".docx":
        return extract_text_from_docx(path)
    if extension == ".pptx":
        return extract_text_from_pptx(path)
    if extension == ".txt":
        try:
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()
        except Exception as ex:
            logger.warning("Failed to read text file '%s': %s", path, ex)
            return ""
    # For images and unrecognized types, fallback to OCR
    return extract_text_from_image(path)
